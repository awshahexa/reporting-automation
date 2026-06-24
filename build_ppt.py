"""Build interactive PPT for the new SharePoint folder management flow."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
BLUE = RGBColor(0x29, 0x75, 0xB5)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
GRAY = RGBColor(0x66, 0x70, 0x85)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
MEDIUM_GRAY = RGBColor(0x9C, 0xA3, 0xAF)


def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, text="", font_size=12, font_color=WHITE, bold=False, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.margin_left = Pt(6)
    shape.text_frame.margin_right = Pt(6)
    shape.text_frame.margin_top = Pt(4)
    shape.text_frame.margin_bottom = Pt(4)
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, font_color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_bullet_box(slide, left, top, width, height, items, font_size=13, font_color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.space_after = Pt(6)
    return txBox


def add_arrow(slide, left, top, width, height, color=GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width, height, color=GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.2),
             "SharePoint Folder Management Flow", font_size=40, font_color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(0.8),
             "Document Control & Archive Pipeline for Reporting Automation", font_size=20, font_color=MEDIUM_GRAY)
add_shape(slide, Inches(0.8), Inches(3.8), Inches(3.5), Inches(0.6), BLUE,
          "Sacofa — Reporting Automation Project", font_size=14, bold=True)
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.5),
             "Working Folders → Archive Folders → Reporting Automation", font_size=14, font_color=GRAY)

# ═══════════════════════════════════════════
# SLIDE 2: Current Problem
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Current Situation", font_size=32, font_color=WHITE, bold=True)

add_bullet_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5), [
    " Documents are submitted, reviewed, and approved in the same SharePoint location",
    " No separation between work-in-progress and final approved documents",
    " Reporting Automation uses hot_folder/ for PDF extraction — no connection to SharePoint",
    " No audit trail for document lifecycle (submitted → reviewed → approved → archived)",
    " Document Controller has no centralized database to track document history",
    " Risk of approved documents being accidentally modified or deleted",
])
for i, s in enumerate(slide.shapes):
    if hasattr(s, 'text') and s.text.startswith(" "):
        s.text_frame.paragraphs[0].font.color.rgb = RED
        s.text_frame.paragraphs[0].font.size = Pt(15)

# ═══════════════════════════════════════════
# SLIDE 3: Solution Overview
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Solution: Two-Zone Folder Architecture", font_size=32, font_color=WHITE, bold=True)

# Zone 1 box
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5), BLUE,
          "WORKING FOLDERS\nDocuments in Progress", font_size=18, bold=True)
add_bullet_box(slide, Inches(1.0), Inches(2.3), Inches(5.0), Inches(1.5), [
    "• Files are submitted, reviewed, and approved",
    "• 3 main areas: PMC / Project / Sites Document",
    "• Each has Submit → Review → Approve stages",
    "• Active collaboration — multiple users",
], font_size=13)

# Arrow
add_arrow(slide, Inches(6.5), Inches(2.3), Inches(1.2), Inches(0.5), GREEN)

# Zone 2 box
add_shape(slide, Inches(8.0), Inches(1.5), Inches(5.0), Inches(2.5), GREEN,
          "ARCHIVE FOLDERS\nFinal Approved Documents", font_size=18, bold=True)
add_bullet_box(slide, Inches(8.2), Inches(2.3), Inches(4.5), Inches(1.5), [
    "• Read-only — single approved version only",
    "• Mirrors Working structure (PMC/Project/Sites Doc)",
    "• Source for Reporting Automation extraction",
    "• Strict access: 1 person + admin",
], font_size=13)

add_bullet_box(slide, Inches(0.8), Inches(4.5), Inches(12), Inches(2.5), [
    "Key Design Principles:",
    "  1. Working = staging area, Archive = source of truth",
    "  2. Automated move script transfers approved docs from Working → Archive",
    "  3. Every archive event logged to Document Controller database",
    "  4. Delta Query tracks all changes in both zones",
    "  5. Reporting Automation consumes from Archive/Sites Document only",
], font_size=14, font_color=MEDIUM_GRAY)

# ═══════════════════════════════════════════
# SLIDE 4: Folder Structure
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Complete Folder Structure", font_size=32, font_color=WHITE, bold=True)

# Working
add_shape(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(0.5), BLUE, "WORKING FOLDERS", bold=True, font_size=14)
y = Inches(1.9)
folders = [
    ("PMC Folder", ["Submit", "Review", "Approve"]),
    ("Project Folder", ["Submit", "Review", "Approve"]),
    ("Sites Document", ["Submit", "Review", "Approve"]),
]
for fname, stages in folders:
    add_shape(slide, Inches(0.7), y, Inches(2.0), Inches(0.4), RGBColor(0x3B, 0x82, 0xF6), fname, font_size=11, bold=True)
    for i, st in enumerate(stages):
        cx = Inches(2.9) + Inches(i * 1.4)
        colors = [ORANGE, RGBColor(0x8B, 0x5C, 0xF6), GREEN]
        add_shape(slide, cx, y, Inches(1.2), Inches(0.4), colors[i], st, font_size=10)
    y += Inches(0.5)

add_text_box(slide, Inches(0.7), y + Inches(0.1), Inches(5.5), Inches(0.4),
             "Each stage has subfolders: MoM, Audio&Video, TSSR, TSS, PO, DN, BL, PL, SATP, MOP, SMR, RFI, Drawing, Report, Correspondence",
             font_size=10, font_color=MEDIUM_GRAY)

# Archive
add_shape(slide, Inches(7.5), Inches(1.3), Inches(5.5), Inches(0.5), GREEN, "ARCHIVE FOLDERS", bold=True, font_size=14)
ay = Inches(1.9)
for fname in ["PMC Folder", "Project Folder", "Sites Document"]:
    add_shape(slide, Inches(7.7), ay, Inches(2.5), Inches(0.4), RGBColor(0x16, 0xA3, 0x4A), fname, font_size=11, bold=True)
    add_shape(slide, Inches(10.5), ay, Inches(2.0), Inches(0.4), GREEN, "Doc Types", font_size=10)
    ay += Inches(0.5)

# Arrow between
add_arrow(slide, Inches(6.2), Inches(2.8), Inches(1.2), Inches(0.4), GREEN)
add_text_box(slide, Inches(6.2), Inches(3.3), Inches(1.2), Inches(0.4),
             "Auto Move", font_size=10, font_color=GREEN, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(7.5), Inches(3.8), Inches(5.5), Inches(1.5),
             "Archive is the Document Control repository:\n• Only 1 approved version per document\n• Read-only for most users\n• Full audit trail via DC database\n• Sites Document/Archive → feeds Reporting Automation",
             font_size=13, font_color=WHITE)

# ═══════════════════════════════════════════
# SLIDE 5: Move + Log Flow
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Automated Archive Pipeline", font_size=32, font_color=WHITE, bold=True)

# Step boxes
steps = [
    ("1. File Approved", "Document moved to\nWorking/*/Approve/{type}/", BLUE),
    ("2. Script Detects", "Watchdog scans Approve\nfolders every 30s", BLUE),
    ("3. Moves to Archive", "File moved to\nArchive/*/{type}/", GREEN),
    ("4. DC Database", "Record created:\ntimestamp, author, checksum", GREEN),
    ("5. Delta Query", "Change tracked:\nadded, modified, deleted", ORANGE),
]

x_start = Inches(0.4)
for i, (title, desc, color) in enumerate(steps):
    x = x_start + Inches(i * 2.5)
    add_shape(slide, x, Inches(1.5), Inches(2.2), Inches(1.4), color, title, font_size=12, bold=True)
    add_shape(slide, x, Inches(3.0), Inches(2.2), Inches(1.0), RGBColor(0x2D, 0x3A, 0x4F), desc, font_size=11)
    if i < len(steps) - 1:
        add_arrow(slide, x + Inches(2.2), Inches(2.0), Inches(0.3), Inches(0.3), MEDIUM_GRAY)

# DC Database detail
add_shape(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5), RGBColor(0x2D, 0x3A, 0x4F), "", font_size=10)
add_text_box(slide, Inches(1.2), Inches(4.6), Inches(5), Inches(0.4),
             "Document Controller (DC) Database — Every Archive Event Captures:", font_size=14, font_color=WHITE, bold=True)

fields = [
    ("Filename", "BETO_TSSR.pdf"),
    ("Doc Type", "TSSR"),
    ("Source", "Working/Sites Document/Approve/TSSR"),
    ("Archive Path", "Archive/Sites Document/TSSR"),
    ("Archived At", "2026-06-11 02:17:56"),
    ("File Size", "0 bytes (SHA256: e3b0c44...)"),
    ("Author", "Extracted from file metadata"),
]
for i, (label, val) in enumerate(fields):
    col = i % 3
    row = i // 3
    x = Inches(1.2) + Inches(col * 3.8)
    y = Inches(5.1) + Inches(row * 0.6)
    add_shape(slide, x, y, Inches(3.5), Inches(0.45), RGBColor(0x1E, 0x29, 0x3B),
              f"{label}: {val}", font_size=10, font_color=MEDIUM_GRAY)

# ═══════════════════════════════════════════
# SLIDE 6: Delta Query
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Change Tracking with Delta Query", font_size=32, font_color=WHITE, bold=True)

add_shape(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(3.5), RGBColor(0x2D, 0x3A, 0x4F), "")
add_bullet_box(slide, Inches(1.2), Inches(1.5), Inches(10.5), Inches(3.0), [
    "What is Delta Query?",
    "  Microsoft Graph API feature that returns only what changed since your last check",
    "",
    "How it works:",
    "  1. First call:  GET /drives/{id}/root/delta  → returns all files + deltaToken",
    "  2. Second call: GET .../delta?token={deltaToken}  → returns only changes",
    "  3. Results include: added files, modified files, and deleted files (@removed)",
    "",
    "What it tracks across all folders (Working + Archive):",
    "  ✅ New files uploaded      ✅ Files modified     ❌ Files deleted",
    "  ✅ File renamed            ✅ Moved between folders",
    "  ✅ Who last modified it    ✅ Timestamp of each change",
], font_size=14)

add_shape(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0), RGBColor(0x2D, 0x3A, 0x4F), "")
add_text_box(slide, Inches(1.2), Inches(5.1), Inches(10.5), Inches(0.4),
             "Example Delta Query Response:", font_size=14, font_color=WHITE, bold=True)
add_text_box(slide, Inches(1.2), Inches(5.6), Inches(10.5), Inches(1.2),
             '{\n  "deltaToken": "abc123...",\n  "added": [ { "name": "BETO_TSSR.pdf", "size": 5242880, "lastModified": "2026-06-11T10:17:56Z" } ],\n  "removed": [ { "name": "old_draft.pdf", "@removed": "deleted" } ],\n  "modified": []\n}',
             font_size=11, font_color=GREEN)

# ═══════════════════════════════════════════
# SLIDE 7: Integration with Reporting Automation
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Integration with Reporting Automation", font_size=32, font_color=WHITE, bold=True)

# Flow diagram
boxes = [
    ("SharePoint\nArchive/Sites Doc", Inches(0.5), Inches(1.5), GREEN),
    ("Delta Query\nWatchdog", Inches(3.5), Inches(1.5), BLUE),
    ("Extraction\nEngine", Inches(6.5), Inches(1.5), ORANGE),
    ("Milestone\nSync", Inches(9.5), Inches(1.5), RGBColor(0x8B, 0x5C, 0xF6)),
]

for text, x, y, color in boxes:
    add_shape(slide, x, y, Inches(2.5), Inches(1.2), color, text, font_size=13, bold=True)

# Arrows
add_arrow(slide, Inches(3.2), Inches(2.0), Inches(0.3), Inches(0.3), MEDIUM_GRAY)
add_arrow(slide, Inches(6.2), Inches(2.0), Inches(0.3), Inches(0.3), MEDIUM_GRAY)
add_arrow(slide, Inches(9.2), Inches(2.0), Inches(0.3), Inches(0.3), MEDIUM_GRAY)

# Dashboard
add_shape(slide, Inches(4.5), Inches(3.5), Inches(4.5), Inches(0.8), GREEN,
          "Dashboard (localhost:8080)", font_size=14, bold=True)
add_down_arrow(slide, Inches(6.5), Inches(2.8), Inches(0.4), Inches(0.5), GREEN)

add_bullet_box(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.5), [
    "End-to-End Flow:",
    "  1. Documents approved in Working → auto-moved to Archive/Sites Document/*",
    "  2. Delta Query detects new files in Archive/Sites Document",
    "  3. New files downloaded and passed to Extraction Engine",
    "  4. Extracted data stored in reporting.db → milestone status updated",
    "  5. Dashboard reflects real-time progress (CC/PAC/FAC completion)",
    "",
    "Difference from current system:",
    "  Before: Manual drop to hot_folder/ → now: Auto-pickup from SharePoint Archive",
    "  Before: No audit trail → now: Full DC database + Delta Query tracking",
    "  Before: Single working folder → now: Two-zone (Working + Archive)",
], font_size=13)

# ═══════════════════════════════════════════
# SLIDE 8: Requirements Checklist
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Requirements — What We Need from Sacofa IT", font_size=32, font_color=WHITE, bold=True)

items = [
    ("Azure App Registration", "Reporting-Automation-App (single tenant)", BLUE),
    ("API Permission", "Sites.ReadWrite.All (application)", BLUE),
    ("Client Secret", "For app-only authentication", BLUE),
    ("Tenant ID", "Sacofa's Microsoft 365 tenant identifier", BLUE),
    ("Client ID", "Application ID from App Registration", BLUE),
    ("SharePoint Site URL", "E.g., https://sacofa.sharepoint.com/sites/Reporting", GREEN),
    ("Document Library", "Library name where folders will be created", GREEN),
    ("Create Folder Structure", "We provide the script — IT runs in SharePoint", GREEN),
    ("Admin Consent", "Global Admin must grant consent for permissions", ORANGE),
]

for i, (label, desc, color) in enumerate(items):
    y = Inches(1.4) + Inches(i * 0.65)
    add_shape(slide, Inches(0.8), y, Inches(3.0), Inches(0.5), color, label, font_size=12, bold=True)
    add_text_box(slide, Inches(4.0), y + Pt(2), Inches(8.5), Inches(0.5), desc, font_size=12, font_color=MEDIUM_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
             "See sharepoint_setup_guide.md for full step-by-step instructions to hand to Sacofa IT.",
             font_size=12, font_color=ORANGE, bold=True)

# ═══════════════════════════════════════════
# SLIDE 9: Commands Reference
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Commands Reference", font_size=32, font_color=WHITE, bold=True)

cmds = [
    ("python run.py sp-setup", "Create folder structure (Working + Archive) — 180 folders"),
    ("python run.py sp-validate", "Check if all expected folders exist"),
    ("python run.py sp-archive", "Move all files from Working/*/Approve to Archive/"),
    ("python run.py sp-watch", "Continuous watchdog — archives files every 30s"),
    ("python run.py sp-delta", "Run simulated Delta Query — show folder state"),
    ("python run.py dc-log", "View Document Controller archive log entries"),
    ("python run.py dc-summary", "Show DC stats: total archived, today, by doc type"),
    ("python run.py start", "Start Reporting Automation dashboard (port 8080)"),
]

for i, (cmd, desc) in enumerate(cmds):
    y = Inches(1.3) + Inches(i * 0.6)
    add_shape(slide, Inches(0.5), y, Inches(5.5), Inches(0.45), BLUE, cmd, font_size=11, bold=True)
    add_text_box(slide, Inches(6.3), y + Pt(2), Inches(6.5), Inches(0.45), desc, font_size=11, font_color=MEDIUM_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
             "All commands work in local simulation mode (sharepoint_sim/) — no SharePoint credentials needed for testing.",
             font_size=12, font_color=ORANGE, bold=True)

# ═══════════════════════════════════════════
# SLIDE 10: Summary
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Summary", font_size=32, font_color=WHITE, bold=True)

points = [
    ("Two-Zone Architecture", "Working folders for in-progress documents\nArchive folders for final approved versions only"),
    ("Automated Move + Log", "Approved documents automatically moved from Working → Archive\nEvery move logged to Document Controller database"),
    ("Delta Query Tracking", "Microsoft Graph API tracks every file change (add/modify/delete)\nNo need to poll or compare snapshots"),
    ("Reporting Automation Ready", "Archive/Sites Document becomes the source for PDF extraction\nDashboard continues to show real-time milestone progress"),
    ("Fully Simulated", "All code works locally without SharePoint credentials\n180 folders created, move+log verified, ready for IT setup"),
]

for i, (title, desc) in enumerate(points):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(1.3) + Inches(row * 2.8)
    colors = [BLUE, GREEN, ORANGE, RGBColor(0x8B, 0x5C, 0xF6), RED]
    add_shape(slide, x, y, Inches(3.8), Inches(0.6), colors[i], title, font_size=14, bold=True)
    add_shape(slide, x, y + Inches(0.7), Inches(3.8), Inches(1.6), RGBColor(0x2D, 0x3A, 0x4F), desc, font_size=11)

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
             "Next step: Provide sharepoint_setup_guide.md to Sacofa IT → get credentials → switch simulate=False → go live",
             font_size=14, font_color=GREEN, bold=True)

# Save
output_path = r"C:\Users\Apex-Guest\Documents\Nizam\Sacofa\NewAddedFlow_Interactive.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
